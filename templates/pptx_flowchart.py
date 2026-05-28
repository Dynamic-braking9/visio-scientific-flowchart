#!/usr/bin/env python3
"""
PowerPoint 科研流程图生成器
使用 python-pptx 在 PPTX 中绘制可编辑的流程图。

依赖: pip install python-pptx

用法:
    from pptx_flowchart import PptxFlowchart
    fc = PptxFlowchart("output.pptx", title="实验流程")
    fc.add_process("数据采集", phase="data")
    fc.save()
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from typing import Optional, Tuple


# ── 配色方案 ────────────────────────────────────

PHASE_COLORS = {
    "data":       (RGBColor(0xDB, 0xEA, 0xFE), RGBColor(0x3B, 0x82, 0xF6)),
    "preprocess": (RGBColor(0xED, 0xE9, 0xFE), RGBColor(0x8B, 0x5C, 0xF6)),
    "analysis":   (RGBColor(0xD1, 0xFA, 0xE5), RGBColor(0x10, 0xB9, 0x81)),
    "check":      (RGBColor(0xFE, 0xF3, 0xC7), RGBColor(0xF5, 0x9E, 0x0B)),
    "verify":     (RGBColor(0xFF, 0xED, 0xD5), RGBColor(0xF9, 0x73, 0x16)),
    "result":     (RGBColor(0xFC, 0xE7, 0xF3), RGBColor(0xEC, 0x48, 0x99)),
}

TERMINAL_COLOR  = (RGBColor(0xF3, 0xF4, 0xF6), RGBColor(0x6B, 0x72, 0x80))
IO_COLOR         = (RGBColor(0xCF, 0xFA, 0xFE), RGBColor(0x06, 0xB6, 0xD4))
SUBPROC_COLOR    = (RGBColor(0xFE, 0xF9, 0xC3), RGBColor(0xCA, 0x8A, 0x04))
TEXT_COLOR        = RGBColor(0x1F, 0x29, 0x37)
ARROW_COLOR       = RGBColor(0x6B, 0x72, 0x80)


class PptxFlowchart:
    """PowerPoint 流程图生成器，生成可编辑的 PPTX 形状。"""

    def __init__(self, filepath: str, title: str = "",
                 direction: str = "vertical",
                 slide_width: float = 13.333, slide_height: float = 7.5):
        """
        Args:
            filepath: 输出 .pptx 文件路径
            title: 流程图标题
            direction: vertical 或 horizontal
            slide_width: 幻灯片宽度（英寸，默认 16:9 宽屏）
            slide_height: 幻灯片高度（英寸）
        """
        self.filepath = filepath
        self.title = title
        self.direction = direction
        self.slide_w = slide_width
        self.slide_h = slide_height

        self.prs = Presentation()
        self.prs.slide_width = Inches(slide_width)
        self.prs.slide_height = Inches(slide_height)

        # 添加空白幻灯片
        layout = self.prs.slide_layouts[6]  # blank layout
        self.slide = self.prs.slides.add_slide(layout)

        self.shapes = []
        self._node_w = 2.5    # inches
        self._node_h = 0.6
        self._h_gap = 0.5
        self._v_gap = 0.4
        self._next_x = 0.0
        self._next_y = 0.0
        self._center_x = slide_width / 2

        # 标题
        if title:
            tx = Inches(0.5)
            ty = Inches(0.2)
            tw = Inches(slide_width - 1)
            th = Inches(0.5)
            title_box = self.slide.shapes.add_textbox(tx, ty, tw, th)
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = TEXT_COLOR
            p.alignment = PP_ALIGN.CENTER
            self._next_y = 1.0
            self._next_x = self._center_x - self._node_w / 2

    def _add_shape(self, text: str, shape_type: str = "process",
                   phase: str = "", w: float = None, h: float = None,
                   x: float = None, y: float = None):
        """添加流程图形状"""
        w = w or self._node_w
        h = h or self._node_h
        x = x if x is not None else self._next_x
        y = y if y is not None else self._next_y

        # 选择形状类型
        shape_map = {
            "process":    MSO_SHAPE.ROUNDED_RECTANGLE,
            "decision":   MSO_SHAPE.DIAMOND,
            "terminal":   MSO_SHAPE.ROUNDED_RECTANGLE,
            "io":         MSO_SHAPE.PARALLELOGRAM,
            "subprocess": MSO_SHAPE.ROUNDED_RECTANGLE,
            "database":   MSO_SHAPE.CAN,
            "annotation": MSO_SHAPE.RIGHT_ARROW_CALLOUT,
        }
        mso_shape = shape_map.get(shape_type, MSO_SHAPE.ROUNDED_RECTANGLE)

        # 获取颜色
        if shape_type == "terminal":
            fill_c, line_c = TERMINAL_COLOR
        elif shape_type == "io":
            fill_c, line_c = IO_COLOR
        elif shape_type == "subprocess":
            fill_c, line_c = SUBPROC_COLOR
        elif phase and phase in PHASE_COLORS:
            fill_c, line_c = PHASE_COLORS[phase]
        else:
            fill_c, line_c = (RGBColor(0xF3, 0xF4, 0xF6), RGBColor(0x9C, 0xA3, 0xAF))

        # 创建形状
        left = Inches(x)
        top = Inches(y)
        width = Inches(w)
        height = Inches(h)

        shape = self.slide.shapes.add_shape(mso_shape, left, top, width, height)

        # 设置填充和边框
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_c
        shape.line.color.rgb = line_c
        shape.line.width = Pt(2)

        # 设置文字
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = TEXT_COLOR
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 子流程双边框效果（添加内矩形）
        if shape_type == "subprocess":
            inner = self.slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + 0.05), Inches(y + 0.05),
                Inches(w - 0.1), Inches(h - 0.1)
            )
            inner.fill.background()
            inner.line.color.rgb = line_c
            inner.line.width = Pt(1)

        self.shapes.append(shape)

        # 更新位置
        if self.direction == "vertical":
            self._next_y = y + h + self._v_gap
        else:
            self._next_x = x + w + self._h_gap

        return shape

    # ── 公开 API ──────────────────────────────

    def add_start(self, text: str = "开始"):
        return self._add_shape(text, "terminal", w=1.8, h=0.5)

    def add_end(self, text: str = "结束"):
        return self._add_shape(text, "terminal", w=1.8, h=0.5)

    def add_process(self, text: str, phase: str = "analysis"):
        return self._add_shape(text, "process", phase=phase)

    def add_decision(self, text: str, phase: str = "check"):
        return self._add_shape(text, "decision", phase=phase, w=2.2, h=1.0)

    def add_io(self, text: str, phase: str = "data"):
        return self._add_shape(text, "io", phase=phase)

    def add_subprocess(self, text: str, phase: str = "analysis"):
        return self._add_shape(text, "subprocess", phase=phase)

    def add_database(self, text: str, phase: str = "data"):
        return self._add_shape(text, "database", phase=phase, h=0.7)

    def add_annotation(self, text: str):
        return self._add_shape(text, "annotation", w=2.0, h=0.45)

    def connect(self, src_idx: int, dst_idx: int, label: str = ""):
        """连接两个形状（按添加顺序的索引）"""
        if src_idx >= len(self.shapes) or dst_idx >= len(self.shapes):
            raise IndexError("形状索引越界")

        src = self.shapes[src_idx]
        dst = self.shapes[dst_idx]

        # 添加连接线（带箭头）
        connector = self.slide.shapes.add_connector(
            1,  # straight connector
            int(src.left + src.width / 2), int(src.top + src.height),
            int(dst.left + dst.width / 2), int(dst.top)
        )
        connector.line.color.rgb = ARROW_COLOR
        connector.line.width = Pt(1.5)

        # 添加箭头
        connector.line.end_marker_style = 2  # triangle arrow

        if label:
            # 添加标签文本框
            lx = int((src.left + dst.left) / 2) + Inches(0.3)
            ly = int((src.top + src.height + dst.top) / 2)
            tb = self.slide.shapes.add_textbox(lx, ly, Inches(0.6), Inches(0.3))
            tf = tb.text_frame
            tf.paragraphs[0].text = label
            tf.paragraphs[0].font.size = Pt(10)
            tf.paragraphs[0].font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

    def auto_connect(self):
        """按添加顺序自动连接所有形状"""
        for i in range(len(self.shapes) - 1):
            self.connect(i, i + 1)

    def save(self):
        """保存 PPTX 文件"""
        self.prs.save(self.filepath)
        print(f"✅ PPTX 流程图已生成: {self.filepath}")


# ── 演示 ─────────────────────────────────────────

def demo():
    fc = PptxFlowchart("示例_深度学习流程.pptx",
                       title="基于深度学习的图像分类研究流程")

    fc.add_start("开始")
    fc.add_io("CIFAR-10 数据集", phase="data")
    fc.add_process("数据增强与预处理", phase="data")
    fc.add_process("ResNet-50 模型构建", phase="analysis")
    fc.add_process("模型训练 (Adam)", phase="analysis")
    fc.add_decision("准确率 ≥ 90%？", phase="check")
    fc.add_subprocess("交叉验证", phase="verify")
    fc.add_io("实验报告", phase="result")
    fc.add_end("结束")

    fc.auto_connect()
    fc.save()


if __name__ == "__main__":
    demo()
